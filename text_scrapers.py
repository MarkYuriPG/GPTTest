import csv
import pandas as pd
from docx import Document
from pptx import Presentation


def scrape_txt(path: str):
    with open(path, 'r', encoding='utf-8') as file:
        text_data = file.read()
    print(text_data)

def scrape_docx(path: str):
    doc = Document(path)
    text_data = ''
    for paragraph in doc.paragraphs:
        text_data += paragraph.text
    print(text_data)

def scrape_pptx(path: str):
    ppt = Presentation(path)
    text_data = ''
    for slide_number, slide in enumerate(ppt.slides):
        text_data += f"\nSlide {slide_number+1}:\n"
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text_data += shape.text + '\n'
    print(text_data)

def scrape_csv(path: str):
    text_data = ''
    with open(path, 'r') as csv_file:
        reader = csv.reader(csv_file)
        for row in reader:
            text_data += ','.join(row)
    print(text_data)

def scrape_xlsx(path: str):
    df = pd.read_excel(path)
    text_data = df.to_string(index = False)
    print(text_data)

#scrape_txt("data_privacy.txt")
#scrape_docx("test.docx")
#scrape_pptx("test.pptx")
#scrape_csv("GroceryStoreDataSet.csv")
#scrape_xlsx("test.xlsx")