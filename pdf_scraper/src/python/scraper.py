import csv
import fitz
import pandas as pd
from docx import Document
from pptx import Presentation

class Scraper:
    def __init__(self):
        self.text_data = ""

    def scrape_txt(self, path: str):
        with open(path, 'r', encoding='utf-8') as file:
            self.text_data = file.read()
        return self.text_data

    def scrape_docx(self, path: str):
        doc = Document(path)
        for paragraph in doc.paragraphs:
            self.text_data += paragraph.text
        return self.text_data


    def scrape_pdf(self, path: str):
        pdf_document = fitz.open(path)

        for page_number in range(pdf_document.page_count):
            page = pdf_document[page_number]
            page_text = page.get_text()
            self.text_data += page_text

        return self.text_data


    def scrape_pptx(self, path: str):
        ppt = Presentation(path)
        for slide_number, slide in enumerate(ppt.slides):
            self.text_data += f"\nSlide {slide_number+1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    self.text_data += shape.text + '\n'
        return self.text_data

    def scrape_csv(self, path: str):
        with open(path, 'r') as csv_file:
            reader = csv.reader(csv_file)
            for row in reader:
                self.text_data += ','.join(row)
        return self.text_data

    def scrape_xlsx(self, path: str):
        df = pd.read_excel(path)
        self.text_data = df.to_string(index = False)
        return self.text_data

#scrape_txt("data_privacy.txt")
#scrape_docx("test.docx")
#scrape_pptx("test.pptx")
#scrape_csv("GroceryStoreDataSet.csv")
#scrape_xlsx("test.xlsx")
#print(scrape_txt("data_privacy.txt"))